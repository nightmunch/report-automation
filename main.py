import psycopg2
from pptx import Presentation
from pptx.util import Inches

import os
from dotenv import load_dotenv

load_dotenv()

# Connect to PostgresSQL DB
conn = psycopg2.connect(
    host=os.getenv("SQL_HOST"),
    database=os.getenv("SQL_DATABASE"),
    user=os.getenv("SQL_USER"),
    password=os.getenv("SQL_PASSWORD")
)

# Create a cursor object to execute SQL queries
cur = conn.cursor()

# Retrieve data from the database
cur.execute("SELECT * FROM kewangan")
data = cur.fetchall()

# Close the cursor and connection to the database
cur.close()
conn.close()

# Create a new PowerPoint presentation
prs = Presentation()

columns = ["Inisiatif", "Agensi", "Perbelanjaan (RM juta)"]

# Iterate over the data and add slides to the presentation
for i in range(0, len(data), 5):
    # Create a new slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Create a textbox on the slide
    textbox = slide.shapes.add_textbox(left=Inches(1), top=Inches(0.5), width=Inches(8), height=Inches(1.5))
    # Create a table on the slide
    table = slide.shapes.add_table(rows=min(6, len(data)-i+1), cols=3, left=Inches(1), top=Inches(1), width=Inches(8), height=Inches(1.5)).table

    # add the text "Menteri Kewangan" to the text box
    textframe = textbox.text_frame
    textframe.text = f"Menteri Kewangan ({int(i/len(data)/5)+1}/{int(len(data)/5)})"

    # Add headers to the table
    for count, column in enumerate(columns):
        table.cell(0, count).text = column
        for j in range(i, min(i+5, len(data))):
            table.cell(j-i+1, count).text = str(data[j][count])

# Save the PowerPoint presentation
prs.save('output/Laporan_Bulanan.pptx')